#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""将 KVCache 管理机制PPT 简化为极简风格：
- 删除无文字的空填充装饰块
- 所有含文字的图形转纯文本框（去填充、去线框、去阴影）
- 文字统一改为近黑色
输出到新文件，保留原文件。
"""
import copy
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree

SRC = r"c:/Users/89517/Desktop/vllm同步/vllm/vllm/v1/core/kvcache_ppt/KVCache_管理机制详解.pptx"
DST = r"c:/Users/89517/Desktop/vllm同步/vllm/vllm/v1/core/kvcache_ppt/KVCache_管理机制详解（极简）.pptx"

INK = RGBColor(0x26, 0x26, 0x26)   # 近黑色正文

ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def has_text(shape):
    if shape.has_text_frame:
        return bool(shape.text_frame.text.strip())
    return False


def strip_shadow(shape):
    """移除阴影/发光等效果层 a:effectLst / a:effectDag"""
    spPr = shape._element.spPr
    if spPr is None:
        return
    for tag in ("a:effectLst", "a:effectDag"):
        for el in spPr.findall(tag, ns):
            spPr.remove(el)


def strip_text_belt_shape(shape):
    """去填充、去线框、去阴影 -> 纯文本框"""
    try:
        shape.fill.background()          # noFill
    except Exception:
        pass
    try:
        shape.line.fill.background()     # 无线框
    except Exception:
        pass
    strip_shadow(shape)


def normalize_text(shape):
    try:
        tf = shape.text_frame
    except Exception:
        return
    for r in tf.paragraphs:
        for run in r.runs:
            try:
                run.font.color.rgb = INK
            except Exception:
                try:
                    run.font.color.theme_color = None
                except Exception:
                    pass
            try:
                run.font.highlight_color = None
            except Exception:
                pass


def process_shape(shape):
    """返回 True 表示该 shape 应被删除"""
    text = has_text(shape)
    if text:
        # 内容（含标题）-> 纯文本框
        try:
            if shape.shape_type in (1,):  # AUTO_SHAPE
                strip_text_belt_shape(shape)
        except Exception:
            pass
        normalize_text(shape)
        return False
    else:
        # 无文字：若为图形一律删除（清理装饰）
        try:
            if shape.shape_type == 1:
                return True
        except Exception:
            pass
        return False


def main():
    prs = Presentation(SRC)
    removed_total = 0
    kept_total = 0
    for si, slide in enumerate(prs.slides, 1):
        to_delete = []
        kept = 0
        for shape in list(slide.shapes):
            do_del = process_shape(shape)
            if do_del:
                to_delete.append(shape)
            else:
                kept += 1
        for shape in to_delete:
            shape._element.getparent().remove(shape._element)
        kept_total += kept
        removed_total += len(to_delete)
        print(f"slide {si}: kept {kept}, removed {len(to_delete)}")

    prs.save(DST)
    print("\n=== done ===")
    print(f"总保留: {kept_total}, 删除装饰: {removed_total}")
    print("输出:", DST)


if __name__ == "__main__":
    main()