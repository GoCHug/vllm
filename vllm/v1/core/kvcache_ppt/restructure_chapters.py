#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""按用户新目录(8章)重构 PPT：
- 概述=原M1(什么是)+M2(为什么)合并 → 删除原第6页 M2 分隔页
- 8 个分隔页改为「第一章..第八章 · 章名」，编号与溪子同步
- 页内左上「模块 0X · 主题」标签改为「第X章 · 主题」
- 底角页码顺排 1..N
"""
from pptx import Presentation
from pptx.util import Emu

PATH = r"c:/Users/89517/Desktop/vllm同步/vllm/vllm/v1/core/kvcache_ppt/KVCache_管理机制详解_新.pptx"

ORD = [None, "第一章", "第二章", "第三章", "第四章", "第五章", "第六章", "第七章", "第八章"]
SPEC_NUM = ["一", "二", "三", "四", "五", "六", "七", "八"]

# 章名（与用户目录一致）
CH_NAME = {
    1: "KV Cache 概述",
    2: "PagedAttention 原理",
    3: "各类 Attention 的 KV",
    4: "管理机制基础",
    5: "KVCache 五层架构总览",
    6: "端到端生命周期",
    7: "机制逐层剖析",
    8: "总结",
}

# 分隔页(原索引) -> 新章号
DIVIDER = {3: 1, 10: 2, 13: 3, 16: 4, 22: 5, 25: 6, 35: 7, 46: 8}

# 内容页(原索引) -> 新章号
LABEL_CH = {
    4: 1, 5: 1, 7: 1, 8: 1, 9: 1,
    11: 2, 12: 2,
    14: 3, 15: 3,
    17: 4, 18: 4, 19: 4, 20: 4, 21: 4,
    23: 5, 24: 5,
    26: 6, 27: 6, 28: 6, 29: 6, 30: 6, 31: 6, 32: 6, 33: 6, 34: 6,
    36: 7, 37: 7, 38: 7, 39: 7, 40: 7, 41: 7, 42: 7, 43: 7, 44: 7, 45: 7,
    47: 8, 48: 8, 49: 8, 50: 8, 51: 8, 52: 8, 53: 8, 54: 8,
}


def set_run_text(shape, text):
    """改写 shape 首个 run 的文本，清除其余 run/段落，保留基准格式。"""
    tf = shape.text_frame
    base = None
    for para in tf.paragraphs:
        for run in para.runs:
            base = run.font
            break
        if base:
            break
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    if base is not None:
        try:
            run.font.name = base.name
        except Exception:
            pass
        if base.size:
            run.font.size = base.size
        if base.bold is not None:
            run.font.bold = base.bold
        try:
            if base.color and base.color.type == 1:
                run.font.color.rgb = base.color.rgb
        except Exception:
            pass


def find_shape(slide, pred):
    for sh in slide.shapes:
        if sh.has_text_frame and pred(sh.text_frame.text):
            return sh
    return None


def find_shape_fs(slide, lo, hi):
    """按主 run 字号范围定位 shape（用于模块名 29pt）"""
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        m = 0
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    m = max(m, run.font.size.pt)
        if lo <= m <= hi:
            return sh
    return None


def main():
    prs = Presentation(PATH)
    slides = list(prs.slides)

    # ---- 1) 分隔页：溪子 + 大号章号 + 章名 ----
    for orig_idx, ch in DIVIDER.items():
        slide = slides[orig_idx - 1]
        kicker = find_shape(slide, lambda t: t.strip().startswith("MODULE"))
        bignum = find_shape(slide, lambda t: t.strip().isdigit())
        name = find_shape_fs(slide, 28, 30)  # 模块名 29pt 的 shape
        if kicker is None or bignum is None or name is None:
            print(f"[warn] 分隔页 {orig_idx} 结构未完全匹配")
            continue
        set_run_text(bignum, ORD[ch])
        from pptx.util import Pt
        w_pt = bignum.width / 12700
        fs = max(30, min(64, int(w_pt / (len(ORD[ch]) * 0.80))))
        for run in bignum.text_frame.paragraphs[0].runs:
            run.font.size = Pt(fs)
        set_run_text(name, CH_NAME[ch])
        set_run_text(kicker, f"CHAPTER 0{ch} · DEEP DIVE")
        print(f"分隔页原slide{orig_idx} -> {ORD[ch]} {CH_NAME[ch]}")

    # ---- 2) 内容页左上标签 模块 0X -> 第X章 ----
    for orig_idx, ch in LABEL_CH.items():
        slide = slides[orig_idx - 1]
        lab = find_shape(slide, lambda t: t.strip().startswith("模块 "))
        if lab is None:
            continue
        txt = lab.text_frame.text.strip()
        # 拆成 "第X章" + 主题词
        tail = txt.split("·", 1)[-1].strip()
        new_txt = f"第{SPEC_NUM[ch-1]}章 · {tail}" if tail else f"第{SPEC_NUM[ch-1]}章"
        set_run_text(lab, new_txt)
        print(f"  内容slide{orig_idx} 标签 -> {new_txt}")

    # ---- 3) 删除原第6页(M2 分隔页) ----
    # 先定位：当前仍是原顺序
    slide6 = slides[5]
    drop_slide(prs, slide6)
    print("已删除 M2 分隔页(原第6页)")

    # ---- 4) 底角页码顺排 ----
    SW = prs.slide_width / 914400
    SH = prs.slide_height / 914400
    for idx, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and sh.top / 914400 > SH - 0.75 and sh.left / 914400 > SW / 2:
                set_run_text(sh, str(idx + 1))
                break
    print("页码已顺排")

    prs.save(PATH)
    print("已保存:", PATH)


def drop_slide(prs, slide):
    """按 0-based 位置删除幻灯片（删除引用关系 + 移除 sldId）。"""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    idx = -1
    for i, sldId in enumerate(list(sldIdLst)):
        rId = sldId.get(qn("r:id"))
        try:
            part = prs.part.related_part(rId)
        except Exception:
            continue
        if part._element is slide._element:
            idx = i
            target_rId = rId
            break
    if idx < 0:
        raise RuntimeError("找不到要删除的 slide")
    target = list(sldIdLst)[idx]
    sldIdLst.remove(target)
    prs.part.drop_rel(target_rId)
    print("已删除 slide idx", idx)


if __name__ == "__main__":
    main()