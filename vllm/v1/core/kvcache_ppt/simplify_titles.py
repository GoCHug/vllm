#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""精简标题与目录（正文不动），覆盖写入极简版文件。"""
from pptx import Presentation
from pptx.util import Pt

PATH = r"c:/Users/89517/Desktop/vllm同步/vllm/vllm/v1/core/kvcache_ppt/KVCache_管理机制详解（极简）.pptx"

# 页内标题精简映射（slide_idx -> 简洁标题）
NEW_TITLES = {
    4:  "K / V 是什么",
    5:  "KV Cache 定义",
    7:  "重算代价：O(n²)",
    8:  "KV 显存账本",
    9:  "连续存储的三大浪费",
    11: "逻辑连续 · 物理离散",
    12: "三条设计 + 一份成绩",
    14: "三种注意力的 KV",
    15: "数据类总览",
    17: "六个积木词",
    18: "灵感：虚拟内存",
    19: "两把钥匙",
    20: "链式哈希",
    21: "整数索引 · 不搬显存",
    23: "五层全景 · 谁持有谁",
    24: "单模型下的五层",
    26: "示例请求 R",
    27: "get_computed_blocks",
    28: "allocate_slots",
    29: "cache_blocks 前向",
    30: "BlockPool 演算",
    31: "free 逆序回收",
    32: "watermark 与抢占",
    33: "调度器主循环",
    34: "整轮生命周期",
    36: "初始化五步",
    37: "张量 · block_id 即行号",
    38: "KVCacheBlock",
    39: "BlockPool 五条铁律",
    40: "FullAttentionManager",
    41: "Copy-on-Write",
    42: "Coordinator 直通车",
    43: "KVCacheManager 门面",
    44: "Full + 滑窗混合",
    45: "GDN + Full 混合",
    47: "八条设计哲学",
    48: "block_size 的推演",
    49: "前缀缓存收益",
    50: "扩展类型速览",
    51: "四个常见误区",
    52: "八问八答",
    54: "九模块 × 十个关键词",
}

TOC_SUBTITLE = "九模块渐进 · 入门到专家"


def block_max_font(shape):
    m = 0
    if not shape.has_text_frame:
        return m
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                m = max(m, run.font.size.pt)
    return m


def set_shape_text(shape, newtext):
    tf = shape.text_frame
    # 记录基准字体属性（取第一个 run）
    base = None
    for para in tf.paragraphs:
        for run in para.runs:
            base = run.font
            break
        if base:
            break
    # 清空所有段落
    for para in list(tf.paragraphs):
        p = para._p
        p.getparent().remove(p)
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    run = p.add_run()
    run.text = newtext
    if base is not None:
        sz = base.size
        b = base.bold
        i = base.italic
        col = base.color
        try:
            run.font.name = base.name
        except Exception:
            pass
        if sz:
            run.font.size = sz
        if b is not None:
            run.font.bold = b
        if i is not None:
            run.font.italic = i
        try:
            if col and col.type is not None:
                if col.type == 1:
                    run.font.color.rgb = col.rgb
        except Exception:
            run.font.color.rgb = None


def main():
    prs = Presentation(PATH)
    slides = list(prs.slides)

    # 1) 精简页内标题
    for idx, t in NEW_TITLES.items():
        slide = slides[idx - 1]
        best = None
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            fs = block_max_font(sh)
            if 19 <= fs <= 23:
                if best is None or fs > best[1]:
                    best = (sh, fs)
        if best is not None:
            set_shape_text(best[0], t)

    # 2) 目录页(第2页)：替换副题 + 删除小字关键词行(7pt) 与冗余长说明(11pt)
    toc = slides[1]
    for sh in list(toc.shapes):
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        fs = block_max_font(sh)
        txt = sh.text_frame.text.strip()
        if txt == '从"为什么"到"怎么做"：一条学习路径':
            set_shape_text(sh, TOC_SUBTITLE)
        elif fs == 7:
            sh._element.getparent().remove(sh._element)
        elif '对应"入门' in txt:
            sh._element.getparent().remove(sh._element)

    prs.save(PATH)
    print("已更新:", PATH)


if __name__ == "__main__":
    main()